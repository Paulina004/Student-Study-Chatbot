"""
Handles loading, extracting, and indexing content from PDF and PPTX files.
Provides functions for document chunking, metadata creation, and vectorstore management.
"""
import fitz
from pptx import Presentation
import re
import os
import streamlit as st
from langchain_community.vectorstores import FAISS
from langchain.schema import Document



def extract_text_from_pptx(pptx_file):
    """
    Extract text from PowerPoint files with structure and content preservation.
    """
    # Load the PowerPoint file
    prs = Presentation(pptx_file)
    text_chunks = []
    # Loop through each slide in the presentation
    for slide_num, slide in enumerate(prs.slides, 1):
        # Initialize slide components
        slide_title = ""
        slide_content = []
        bullet_points = []
        other_text = []
        # Extract title if it exists
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_title = slide.shapes.title.text.strip() 
        # Process each shape in the slide
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue      
            text = shape.text.strip()
            if not text:
                continue 
            # Skip if this is the title we already processed
            if shape == slide.shapes.title:
                continue     
            # Clean the text while preserving some structure
            lines = text.split('\n')
            cleaned_lines = [] 
            for line in lines:
                line = line.strip()     
                # Clean up the line
                line = re.sub(r'\s+', ' ', line)
                # Check if it's likely a bullet point
                if line.startswith(('•', '-', '*', '○', '·')) or re.match(r'^\d+[\).]', line):
                    bullet_points.append(line.lstrip('•-*○·').strip())
                else:
                    cleaned_lines.append(line) 
            if cleaned_lines:
                other_text.extend(cleaned_lines)
        
        # Construct the final content with clear structure
        content_parts = []
        # Add title if it exists
        if slide_title:
            content_parts.append(f"**Title:** {slide_title}")
        # Add bullet points with structure
        if bullet_points:
            content_parts.append("**Key Points:**")
            content_parts.extend(f"• {point}" for point in bullet_points)
        # Add other text
        if other_text:
            content_parts.append("**Additional Content:**")
            content_parts.extend(other_text)
        if content_parts:  # Only create a chunk if we have content
            # Join all content with clear section separation
            full_content = "\n".join(content_parts)
            # Create metadata for source tracking
            metadata = {
                "source": getattr(pptx_file, 'name', 'Unknown Source'),
                "slide_number": slide_num,
                "slide_title": slide_title or "Untitled",
                "type": "slide",
                "content_type": "presentation",
                "has_bullets": bool(bullet_points),
                "content_sections": len(content_parts)
            }
            # Create Document object with metadata
            doc = Document(
                page_content=full_content,
                metadata=metadata
            )
            text_chunks.append(doc)
    return text_chunks



def extract_text_from_pdf(pdf_file):
    """
    Extract text and preserve structure from PDF files, including tables, lists, and code blocks where possible.
    """
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    text_chunks = []
    for page_num, page in enumerate(doc, 1):
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if block["type"] == 0:  # text block
                lines = block["lines"]
                chunk_lines = []
                for line in lines:
                    spans = line["spans"]
                    for span in spans:
                        text = span["text"].strip()
                        if not text:
                            continue
                        # Try to detect code blocks (monospace font)
                        if "Mono" in span.get("font", ""):
                            text = f"```\n{text}\n```"
                        chunk_lines.append(text)
                if chunk_lines:
                    para = " ".join(chunk_lines)
                    # Try to detect lists (lines starting with bullets or numbers)
                    if any(l.strip().startswith(("-", "•", "*", "·", "○", "1.", "a.", "i.")) for l in chunk_lines):
                        para = "\n".join(chunk_lines)
                    metadata = {
                        "source": getattr(pdf_file, 'name', 'Unknown Source'),
                        "page_number": page_num,
                        "type": "page",
                        "content_type": "pdf"
                    }
                    doc_obj = Document(
                        page_content=para,
                        metadata=metadata
                    )
                    text_chunks.append(doc_obj)
            elif block["type"] == 1:  # image block
                # Insert a placeholder for images
                metadata = {
                    "source": getattr(pdf_file, 'name', 'Unknown Source'),
                    "page_number": page_num,
                    "type": "image",
                    "content_type": "pdf"
                }
                doc_obj = Document(
                    page_content="[Image omitted]",
                    metadata=metadata
                )
                text_chunks.append(doc_obj)
            elif block["type"] == 2:  # table block (if detected)
                # PyMuPDF does not natively extract tables, but we can try to preserve structure
                table_text = block.get("text", "").strip()
                if table_text:
                    metadata = {
                        "source": getattr(pdf_file, 'name', 'Unknown Source'),
                        "page_number": page_num,
                        "type": "table",
                        "content_type": "pdf"
                    }
                    doc_obj = Document(
                        page_content=f"[Table]\n{table_text}",
                        metadata=metadata
                    )
                    text_chunks.append(doc_obj)
    return text_chunks



def load_document_and_index(file, embeddings, faiss_path="vector_index.faiss"):
    """
    Load and index a document with better error handling and file operations.
    """
    try:
        if not hasattr(file, 'name'):
            st.error("Invalid File Object: Missing Filename")
            return None
        # Get directory path
        faiss_dir = os.path.dirname(faiss_path)
        if faiss_dir and not os.path.exists(faiss_dir):
            try:
                os.makedirs(faiss_dir, mode=0o755, exist_ok=True)
            except Exception as e:
                st.error(f"Error Creating Directory {faiss_dir}: {str(e)}")
                # Fall back to temp directory
                import tempfile
                faiss_path = os.path.join(tempfile.gettempdir(), os.path.basename(faiss_path))
        # Extract text based on file type
        if file.name.endswith('.pdf'):
            chunks = extract_text_from_pdf(file)
        elif file.name.endswith(('.pptx', '.ppt')):
            chunks = extract_text_from_pptx(file)
        else:
            st.error(f"Unsupported File Type: {file.name}")
            return None
        if not chunks:
            st.warning(f"No Content Extracted from {file.name}")
            return None
        # Extract text content and metadata
        texts = [doc.page_content for doc in chunks]
        metadatas = [doc.metadata for doc in chunks]
        # Verify metadata structure
        for metadata in metadatas:
            if not isinstance(metadata, dict):
                st.error("Invalid Metadata Structure Detected")
                return None
            if "type" not in metadata:
                metadata["type"] = "unknown"

        try:
            # If the vector store already exists, overwrite it
            if os.path.exists(faiss_path):
                try:
                    os.remove(faiss_path)
                except OSError:
                    pass  # Ignore and overwrite
                if os.path.exists(faiss_path + ".pkl"):
                    try:
                        os.remove(faiss_path + ".pkl")
                    except OSError:
                        pass  # Ignore and overwrite
            # Do NOT generate a new unique filename; always use the same faiss_path
            vectorstore = FAISS.from_texts(texts, embeddings, metadatas=metadatas)
            # Try to save the vectorstore
            try:
                vectorstore.save_local(faiss_path)
            except OSError as e:
                st.warning(f"Could not save vector store to disk: {str(e)}")
                st.info("Continuing with in-memory vector store")
            return vectorstore
        except Exception as e:
            st.error(f"Error creating vector store: {str(e)}")
            return None
    except Exception as e:
        st.error(f"Error processing document: {str(e)}")
        return None



def clear_vector_store(faiss_path="vector_index.faiss"):
    """
    Safely clear the vector store files with error handling.
    """
    success = True
    try:
        # Get the directory path
        faiss_dir = os.path.dirname(faiss_path)
        # Try to remove the main index file
        if os.path.exists(faiss_path):
            try:
                os.remove(faiss_path)
            except OSError as e:
                st.warning(f"Could Not Remove Index File: {str(e)}")
                success = False
        # Try to remove the pickle file
        pkl_path = faiss_path + ".pkl"
        if os.path.exists(pkl_path):
            try:
                os.remove(pkl_path)
            except OSError as e:
                st.warning(f"Could Not Remove Pickle File: {str(e)}")
                success = False
        # Try to clean up any temporary files
        if faiss_dir and os.path.exists(faiss_dir):
            try:
                for f in os.listdir(faiss_dir):
                    if f.startswith("vector_index_") and (f.endswith(".faiss") or f.endswith(".pkl")):
                        try:
                            os.remove(os.path.join(faiss_dir, f))
                        except OSError:
                            success = False
            except OSError:
                success = False
        return success
    except Exception as e:
        st.error(f"Error during Cleanup: {str(e)}")
        return False 